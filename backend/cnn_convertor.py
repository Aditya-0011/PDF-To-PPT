import keras_ocr
import cv2
import math
import numpy as np
from PIL import Image
import fitz
import io
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pdf2docx import Converter

def pdf_to_images(pdf_path, output_folder='pdf_images'):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    doc = fitz.open(pdf_path)
    image_paths = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap()
        image_path = os.path.join(output_folder, f'page_{page_num}.png')
        pix.save(image_path)
        image_paths.append(image_path)

    doc.close()
    return image_paths


def pdf_to_ppt(pdf_path, output_ppt_path=None):
    if output_ppt_path is None:
        output_ppt_path = f'{Path().absolute()}\\pptx\\{str(dir_name)}.pptx'

    image_paths = pdf_to_images(pdf_path)
    prs = Presentation()

    all_detected_text = []

    for image_path in image_paths:
        detected_text = remove_text_and_create_ppt(image_path, prs)
        all_detected_text.append(detected_text)

    prs.save(output_ppt_path)

    return ' '.join(all_detected_text), output_ppt_path


def image_to_ppt(image_path, output_ppt_path=None):
    if output_ppt_path is None:
        output_ppt_path = os.path.splitext(image_path)[0] + '.pptx'

    prs = Presentation()
    detected_text = remove_text_and_create_ppt(image_path, prs)
    prs.save(output_ppt_path)

    return detected_text, output_ppt_path


def pixels_to_emus(pixel, dpi=96):
    inches = pixel / dpi
    return Emu(inches * 914400)


def estimate_font_size(height, scaling_factor=0.75):
    return Pt(height * scaling_factor)


def midpoint(x1, y1, x2, y2):
    return (int((x1 + x2) / 2), int((y1 + y2) / 2))


def remove_text_and_create_ppt(image_path, prs):
    pipeline = keras_ocr.pipeline.Pipeline()

    image = keras_ocr.tools.read(image_path)
    h, w, _ = image.shape

    prediction_groups = pipeline.recognize([image])

    mask = np.zeros(image.shape[:2], dtype="uint8")
    text_positions = []
    all_detected_text = []
    for box in prediction_groups[0]:
        text, box_points = box
        all_detected_text.append(text)

        x0, y0, x1, y1, x2, y2, x3, y3 = np.array(box_points).flatten()

        x_mid0, y_mid0 = midpoint(x1, y1, x2, y2)
        x_mid1, y_mid1 = midpoint(x0, y0, x3, y3)

        thickness = int(np.sqrt((x2 - x1) ** 2 + (y2 - y1) ** 2))

        cv2.line(mask, (x_mid0, y_mid0), (x_mid1, y_mid1), 255, thickness)
        text_positions.append((text, x0, y0, x2 - x0, y2 - y0))

    result = cv2.inpaint(image, mask, 7, cv2.INPAINT_NS)
    output_image_path = os.path.splitext(image_path)[0] + '_output.jpg'
    cv2.imwrite(output_image_path, cv2.cvtColor(result, cv2.COLOR_BGR2RGB))

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(output_image_path, Emu(0), Emu(0), width=pixels_to_emus(w), height=pixels_to_emus(h))

    for text, x, y, width, height in text_positions:
        textbox = slide.shapes.add_textbox(pixels_to_emus(x), pixels_to_emus(y),
                                           pixels_to_emus(width), pixels_to_emus(height))
        p = textbox.text_frame.paragraphs[0]
        p.text = text
        p.font.size = estimate_font_size(height)

    return ' '.join(all_detected_text)
dir_name = sys.argv[1]
pptx_dir_path = f'{Path().absolute()}\\pptx\\{str(dir_name)}\\{str(dir_name)}.pptx'
print(pptx_dir_path)
file_path = f'{Path().absolute()}\\pdf\\{str(dir_name)}.pdf'
pdf_to_ppt(file_path,pptx_dir_path)